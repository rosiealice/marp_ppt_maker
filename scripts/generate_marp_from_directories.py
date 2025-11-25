#!/usr/bin/env python3
"""
generate_marp_from_directories.py

Searches multiple directories for PNG figures matching variable names, selects the latest year for each variable, and generates a Marp markdown file for PowerPoint export.

Usage:
    python generate_marp_from_directories.py --dirs dir1 dir2 ... --vars var1 var2 ... [--output slides.md] [--pptx]

Depends on: requests, beautifulsoup4
"""
import sys
import os
import re
import argparse
import subprocess
from glob import glob
from urllib.parse import urlparse, urljoin
import requests
from bs4 import BeautifulSoup


# Helper to extract year from filename

YEAR_RE = re.compile(r'(\d{4})(?=\.png$)')

# Helper to list PNGs in a remote HTML directory
def list_remote_pngs(url):
    try:
        r = requests.get(url, timeout=30)
        r.raise_for_status()
    except Exception as e:
        print(f"Failed to fetch {url}: {e}", file=sys.stderr)
        return []
    soup = BeautifulSoup(r.text, "html.parser")
    pngs = []
    for a in soup.find_all('a', href=True):
        href = a['href']
        if href.lower().endswith('.png'):
            full_url = urljoin(url, href)
            pngs.append(full_url)
    for img in soup.find_all('img', src=True):
        src = img['src']
        if src.lower().endswith('.png'):
            full_url = urljoin(url, src)
            pngs.append(full_url)
    return pngs

# Helper to list PNGs in a local directory
def list_local_pngs(path):
    return [os.path.abspath(p) for p in glob(os.path.join(path, '*.png'))]

# Find best PNGs for each directory and variable
def find_best_pngs_per_dir_var(dirs, vars):
    result = []  # List of (var, dir, png_path)
    for var in vars:
        for d in dirs:
            if d.startswith('http://') or d.startswith('https://'):
                pngs = list_remote_pngs(d)
            else:
                pngs = list_local_pngs(d)
            candidates = []
            for p in pngs:
                fname = os.path.basename(urlparse(p).path)
                if var in fname:
                    m = YEAR_RE.search(fname)
                    year = int(m.group(1)) if m else -1
                    candidates.append((year, p))
            if candidates:
                candidates.sort(reverse=True)
                result.append((var, d, candidates[0][1]))
            else:
                print(f"Warning: No PNG found for variable '{var}' in directory '{d}'", file=sys.stderr)
    return result

def write_marp_md_per_dir_var(var_dir_png_list, out_md, dir_to_comment):
    with open(out_md, 'w', encoding='utf-8') as f:
        f.write("---\n")
        f.write("marp: true\n")
        f.write("theme: default\n")
        f.write("paginate: true\n")
        f.write("---\n\n")
        f.write("# Variable Figures by Directory\n\n")
        f.write("<!-- Auto-generated Marp slides for variable figures by directory -->\n\n")
        for var, d, png_path in var_dir_png_list:
            f.write("---\n")
            comment = dir_to_comment[d] if dir_to_comment and d in dir_to_comment else (d.rstrip('/').split('/')[-2] if d.endswith('/') else d.split('/')[-1])
            f.write(f"### <span style='font-size:0.84em'>{var}</span>\n")
            f.write(f"#### <span style='font-size:0.7em'>{comment}</span>\n\n")
            f.write(f"<div style=\"display:flex;align-items:center;justify-content:center;height:100%;\">\n")
            f.write(f"  <img src=\"{png_path}\" style=\"max-width:100%;max-height:100%;\" alt=\"{var}\">\n")
            f.write("</div>\n\n")
    print(f"Wrote Marp markdown to {out_md}")

def export_pptx_from_marp(md_file, pptx_file):
    print('mdfile',md_file,pptx_file)

    from pptx import Presentation
    from pptx.util import Inches
    from bs4 import BeautifulSoup
    import markdown

# Read your Markdown file  
    with open(md_file, "r", encoding="utf-8") as f:
        md_text = f.read()

# Convert Markdown to HTML
    html = markdown.markdown(md_text, extensions=["extra"])
    
    soup = BeautifulSoup(html, "html.parser")


    prs = Presentation()
    prs.slide_width = Inches(15)  # for a 16:9 ratio  
    prs.slide_height = Inches(5)

    current_slide = None
    text_frame = None

    for element in soup.find_all(['h3', 'h4', 'div']):
        if element.name == 'h3':
        # Create new slide for each H3
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        # Add title from span text
            span = element.find('span')
            if span:
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
                title_box.text = span.get_text()
            current_slide = slide

        elif element.name == 'h4':
            # Add subtitle
            if current_slide:
                span = element.find('span')
                if span:
                    subtitle_box = current_slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.5))
                    subtitle_box.text = span.get_text()

        elif element.name == 'div':
            # Add image if present
            img = element.find('img')
            if img and current_slide:
                img_url = img['src']
            # Download image temporarily
                import requests
                from io import BytesIO
                response = requests.get(img_url)
                image_stream = BytesIO(response.content)
                # Add image to slide
                slide.shapes.add_picture(image_stream, Inches(0), Inches(0.9), width=Inches(15))

# Save PP

    prs.save(pptx_file)
    print("Saved as " ,pptx_file)




    

def main():
    parser = argparse.ArgumentParser(description="Create a Marp presentation from directories of PNG figures for variables.")
    parser.add_argument('--dirs', nargs='+', required=True, help='List of directories to search')
    parser.add_argument('--vars', nargs='+', required=True, help='List of variable names to plot')
    parser.add_argument('--comments', nargs='+', required=True, help='List of comments for each directory (same order as --dirs)')
    parser.add_argument('--output', '-o', default='slides.md', help='Output Markdown filename (default: slides.md)')
    parser.add_argument('--pptx', nargs='?', const='slides.pptx', help='Export PPTX using marp CLI (default: slides.pptx)')
    args = parser.parse_args()

    var_dir_png_list = find_best_pngs_per_dir_var(args.dirs, args.vars)
    if not var_dir_png_list:
        print("No figures found for any variables in any directory.", file=sys.stderr)
        sys.exit(1)

    dir_to_comment = {d: c for d, c in zip(args.dirs, args.comments)}
    write_marp_md_per_dir_var(var_dir_png_list, args.output, dir_to_comment)

    if args.pptx:
        pptx_out = args.pptx if isinstance(args.pptx, str) else 'slides.pptx'
        export_pptx_from_marp(args.output, pptx_out)

if __name__ == '__main__':
    main()
